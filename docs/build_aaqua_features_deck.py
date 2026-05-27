"""Generate a standalone .pptx summarizing AAQUA's test-automation features.
Plain template (no branding) — intended for copying slides into the branded deck.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ACCENT = RGBColor(0x6D, 0x28, 0xD9)   # purple
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_title_slide(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    # accent bar
    bar = s.shapes.add_shape(1, Inches(0), Inches(2.7), Inches(13.333), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.9), Inches(11.7), Inches(1.6))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = subtitle
    r2.font.size = Pt(18); r2.font.color.rgb = MUTED
    return s


def add_feature_slide(title, tagline, bullets, notes=None):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, Inches(0.8), Inches(0.7), Inches(0.18), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    # title
    tb = s.shapes.add_textbox(Inches(1.15), Inches(0.6), Inches(11.4), Inches(1.1))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; r = p.add_run(); r.text = title
    r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = DARK
    p2 = tf.add_paragraph(); r2 = p2.add_run(); r2.text = tagline
    r2.font.size = Pt(15); r2.font.italic = True; r2.font.color.rgb = ACCENT
    # bullets
    body = s.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.5), Inches(5.0))
    bf = body.text_frame; bf.word_wrap = True
    for i, (lead, rest) in enumerate(bullets):
        p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
        p.space_after = Pt(8)
        rb = p.add_run(); rb.text = "▸ " + lead
        rb.font.size = Pt(16); rb.font.bold = True; rb.font.color.rgb = DARK
        if rest:
            rr = p.add_run(); rr.text = " — " + rest
            rr.font.size = Pt(15); rr.font.color.rgb = MUTED
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


add_title_slide(
    "AAQUA — AI-Driven Test Automation",
    "New & enhanced capabilities  •  AI-assisted QA utility platform",
)

add_feature_slide(
    "Project-Scoped Workspaces",
    "Every activity is organized under a project bound to one target application",
    [
        ("Active-project gating", "all services require a selected project; selection persists across reloads"),
        ("One project = one target app", "keeps unrelated targets from being mixed in a single workspace"),
        ("Scope-aware, not restrictive", "non-blocking warning when a URL is outside the bound app — third-party integrations still allowed"),
        ("Results tied to the project", "scans, runs, localization & accessibility results are associated for release readiness"),
    ],
    notes="Project scoping gives every team a clean, app-bound workspace while still allowing legitimate off-target URLs (e.g. third-party integrations) with a gentle warning.",
)

add_feature_slide(
    "Functional Test Generator",
    "Requirements → comprehensive, review-ready functional test cases",
    [
        ("Full coverage", "positive, negative, boundary, edge, security — plus navigation, read-only-field validation, and cancel/discard flows"),
        ("Reviewer-friendly detail", "numbered, self-explanatory steps with preconditions and concrete test data"),
        ("One-click export", "structured “Test Cases” Excel sheet ready for manual execution or review"),
        ("Transparency", "shows how long generation took; resilient to large outputs"),
    ],
    notes="The generator now covers the UI behaviors reviewers care about (cancel, read-only, navigation) and writes steps a first-time reviewer can follow without prior context.",
)

add_feature_slide(
    "API Test Generator",
    "OpenAPI/Swagger, raw spec, file upload, or manual endpoints",
    [
        ("Spec-grounded cases", "assertions and status codes come from the documented API — no invented transport-level negatives"),
        ("Manual-testing ready", "each case includes preconditions and plain-language steps (e.g. for Postman)"),
        ("Process-flow (BPMN) mode", "infers ordered multi-step flows for orchestrated back-ends"),
        ("Runnable output", "exports REST Assured (Java) or Playwright (TS) projects, not just cases"),
    ],
    notes="Two modes: per-endpoint test cases and process-flow chains. Output is grounded in the spec and exportable both as a manual checklist and as runnable automation projects.",
)

add_feature_slide(
    "Framework Generator",
    "Scaffold a production-ready automation framework in seconds",
    [
        ("Multi-stack", "Playwright, Cypress, or Selenium — TypeScript/JavaScript/Java"),
        ("Best-practice structure", "Page Object Model, reporting (Allure/HTML), CI/CD, Docker, parallel execution, logging"),
        ("Self-provisioning", "generated Playwright projects auto-install browsers on npm install — runnable out of the box"),
        ("Zero boilerplate", "download a ready-to-run zip with sample tests and config"),
    ],
    notes="Generates a complete, opinionated framework so teams skip setup. Playwright projects now download their browsers automatically, removing the most common 'tests won't run' issue.",
)

add_feature_slide(
    "Reliability & Scale",
    "Engineered for consistent output and multi-user use",
    [
        ("Consistent AI output", "tuned local-LLM reasoning so detailed generations complete instead of timing out"),
        ("Resilient parsing", "tolerant of large/partial model responses — keeps valid results"),
        ("Built-in security", "SSRF-protected scan targets; Keycloak (OIDC) authentication; role-based access"),
        ("Throughput-aware", "guidance for concurrency/queuing as usage grows across teams"),
    ],
    notes="Reliability work ensures the local LLM returns complete, parseable output under load, with security (auth + SSRF protection) and a path to scale for concurrent users.",
)

out = r"D:\AITesting\docs\AAQUA-Features.pptx"
prs.save(out)
print("Saved:", out, "with", len(prs.slides._sldIdLst), "slides")
