"""Generate a comprehensive 'What is AAQUA' overview deck (.pptx).
Plain template (no branding) for copying slides into the branded deck.
Content is grounded in the AAQUA codebase/architecture docs.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

ACCENT = RGBColor(0x6D, 0x28, 0xD9)
ACCENT2 = RGBColor(0x3B, 0x82, 0xF6)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def title_slide(title, subtitle):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, Inches(0), Inches(2.7), Inches(13.333), Inches(0.08))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(0.8), Inches(2.85), Inches(11.7), Inches(1.8))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = DARK
    p = tf.add_paragraph(); r = p.add_run(); r.text = subtitle
    r.font.size = Pt(18); r.font.color.rgb = MUTED


def content_slide(title, tagline, bullets, notes=None, two_col=False):
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, Inches(0.8), Inches(0.65), Inches(0.18), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.15), Inches(0.55), Inches(11.4), Inches(1.15))
    tf = tb.text_frame; tf.word_wrap = True
    r = tf.paragraphs[0].add_run(); r.text = title
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = DARK
    if tagline:
        p = tf.add_paragraph(); r = p.add_run(); r.text = tagline
        r.font.size = Pt(14); r.font.italic = True; r.font.color.rgb = ACCENT

    if two_col:
        half = (len(bullets) + 1) // 2
        cols = [(Inches(1.0), bullets[:half]), (Inches(7.0), bullets[half:])]
    else:
        cols = [(Inches(1.0), bullets)]

    for left, items in cols:
        body = s.shapes.add_textbox(left, Inches(1.95), Inches(5.9) if two_col else Inches(11.5), Inches(5.1))
        bf = body.text_frame; bf.word_wrap = True
        for i, item in enumerate(items):
            lead, rest = item if isinstance(item, tuple) else (item, "")
            p = bf.paragraphs[0] if i == 0 else bf.add_paragraph()
            p.space_after = Pt(7)
            rb = p.add_run(); rb.text = "▸ " + lead
            rb.font.size = Pt(15); rb.font.bold = True; rb.font.color.rgb = DARK
            if rest:
                rr = p.add_run(); rr.text = " — " + rest
                rr.font.size = Pt(14); rr.font.color.rgb = MUTED
    if notes:
        s.notes_slide.notes_text_frame.text = notes


def arch_slide():
    s = prs.slides.add_slide(BLANK)
    bar = s.shapes.add_shape(1, Inches(0.8), Inches(0.65), Inches(0.18), Inches(0.7))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background()
    tb = s.shapes.add_textbox(Inches(1.15), Inches(0.55), Inches(11.4), Inches(0.8))
    r = tb.text_frame.paragraphs[0].add_run(); r.text = "Architecture"
    r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = DARK

    def box(x, y, w, h, text, color):
        sh = s.shapes.add_shape(5, x, y, w, h)  # rounded rect
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.color.rgb = color
        tf = sh.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = 2
        run = p.add_run(); run.text = text
        run.font.size = Pt(13); run.font.bold = True; run.font.color.rgb = RGBColor(255, 255, 255)
        return sh

    box(Inches(0.7), Inches(2.4), Inches(2.7), Inches(1.1), "React + Vite SPA\n(browser UI)", ACCENT2)
    box(Inches(4.0), Inches(2.4), Inches(2.7), Inches(1.1), "Express API\n(:3001)", ACCENT)
    # backends
    box(Inches(7.4), Inches(1.0), Inches(5.1), Inches(0.85), "Local LLM  •  gpt-oss-20b (on-prem, OpenAI-compatible)", DARK)
    box(Inches(7.4), Inches(2.05), Inches(5.1), Inches(0.85), "OWASP ZAP  •  security scanning", DARK)
    box(Inches(7.4), Inches(3.1), Inches(5.1), Inches(0.85), "Playwright  •  browser automation & test runs", DARK)
    box(Inches(7.4), Inches(4.15), Inches(5.1), Inches(0.85), "PostgreSQL  •  projects, results, governance", DARK)
    box(Inches(4.0), Inches(4.6), Inches(2.7), Inches(0.9), "Keycloak (OIDC)\nauth & roles", ACCENT2)

    note = s.shapes.add_textbox(Inches(0.7), Inches(5.9), Inches(12.0), Inches(1.2))
    nf = note.text_frame; nf.word_wrap = True
    r = nf.paragraphs[0].add_run()
    r.text = ("Browser SPA → Express API → fans out to the on-prem LLM, OWASP ZAP, Playwright and PostgreSQL. "
              "Identity is delegated to Keycloak (OIDC, code+PKCE); the LLM runs inside the network so data never leaves Aaseya.")
    r.font.size = Pt(13); r.font.color.rgb = MUTED
    s.notes_slide.notes_text_frame.text = ("AAQUA is a two-process web app (React frontend + Express backend) that orchestrates four engines: "
                                           "a local LLM for generation/analysis, OWASP ZAP for security, Playwright for execution, and PostgreSQL for persistence.")


# 1. Title
title_slide("AAQUA", "AI-Assisted QA Utility Application  —  generate, run, secure & ship quality software")

# 2. What is AAQUA
content_slide(
    "What is AAQUA?",
    "An AI-driven platform that covers the QA lifecycle in one place",
    [
        ("AI-assisted QA platform", "turns requirements and live apps into test artifacts, runs them, and measures release readiness"),
        ("One workspace, many disciplines", "functional, API, security, localization and accessibility testing together"),
        ("On-prem AI", "powered by a local LLM (gpt-oss) — data stays inside Aaseya's network"),
        ("Web-based", "a React + Express application; nothing to install for end users"),
    ],
    notes="AAQUA = AI-Assisted QA Utility Application. The pitch: one AI-powered web platform that spans the whole QA lifecycle, running on an on-prem LLM so data stays internal.",
)

# 3. Why AAQUA
content_slide(
    "Why AAQUA",
    "The problems it solves",
    [
        ("Slow, inconsistent manual test writing", "AI generates comprehensive, review-ready cases in seconds"),
        ("Fragmented, disconnected tools", "a single platform across functional / API / security / a11y / localization"),
        ("No clear view of release quality", "governance scoring and a release-readiness dashboard"),
        ("Data-privacy concerns with cloud AI", "the LLM runs on-premises — nothing sent to external providers"),
        ("Framework setup overhead", "generates ready-to-run automation projects, not just test text"),
    ],
    notes="Position AAQUA against the real pain: manual effort, tool sprawl, no single quality signal, and cloud-AI data concerns.",
)

# 4. The suite
content_slide(
    "The AAQUA Suite",
    "A toolbox of AI-powered QA services",
    [
        ("Functional Test Generator", "requirements → detailed test cases"),
        ("Test Plan & Test Data Generators", "plans and realistic data sets"),
        ("API Test Generator", "specs/endpoints → cases, flows & code"),
        ("Smart Locators", "stable, AI-improved element locators"),
        ("Framework Generator", "Playwright/Cypress/Selenium scaffolds"),
        ("Migration Service", "convert legacy test projects"),
        ("Test Runner", "execute & auto-heal Playwright suites"),
        ("Localization Tester", "i18n / translation checks"),
        ("Accessibility Scanner", "WCAG / axe-based audits"),
        ("Security Scanner", "OWASP ZAP-driven scanning"),
        ("Release Intelligence", "governance & readiness scoring"),
    ],
    two_col=True,
    notes="This is the breadth slide — AAQUA is a suite, not a single tool. Each item is a dedicated module in the app.",
)

# 5. Architecture
arch_slide()

# 6. Functional generation
content_slide(
    "Test Case & Plan Generation",
    "From a requirement to a review-ready test suite",
    [
        ("Comprehensive coverage", "positive, negative, boundary, edge, security, navigation, read-only-field & cancel flows"),
        ("Reviewer-friendly", "numbered steps with preconditions and concrete test data"),
        ("Test plans & data", "generate structured test plans and realistic test data sets"),
        ("Export", "one-click Excel / JSON for execution, sharing or import"),
    ],
    notes="The functional side: test cases, test plans and test data — all AI-generated and export-ready.",
)

# 7. API testing
content_slide(
    "API Test Generator",
    "OpenAPI/Swagger, raw spec, file upload, or manual endpoints",
    [
        ("Spec-grounded cases", "assertions & status codes derived from the documented API"),
        ("Process-flow (BPMN) mode", "ordered, multi-step flows for orchestrated back-ends"),
        ("Manual + automated", "preconditions & plain-language steps, plus runnable REST Assured / Playwright projects"),
        ("Persona-aware auth", "handles secured endpoints and role-based flows"),
    ],
    notes="API testing supports both per-endpoint cases and multi-step process flows, and emits both manual checklists and runnable automation.",
)

# 8. Frameworks, locators, migration
content_slide(
    "Frameworks, Locators & Migration",
    "Bootstrap and modernize automation assets",
    [
        ("Framework Generator", "Playwright / Cypress / Selenium with POM, reporting, CI/CD, Docker — self-provisioning & runnable out of the box"),
        ("Smart Locators", "generates stable element locators and AI-improves weak ones"),
        ("Migration Service", "converts existing test projects into a target framework"),
    ],
    notes="These accelerate getting teams onto solid automation: scaffold new frameworks, get robust locators, and migrate legacy suites.",
)

# 9. Runner
content_slide(
    "Test Runner & Self-Healing",
    "Execute suites and recover from breakages",
    [
        ("Run Playwright projects", "local path or uploaded ZIP, headless or headed"),
        ("Live log streaming", "watch progress in real time; full logs retained"),
        ("AI auto-heal", "suggests & applies fixes for broken locators/steps"),
        ("Results feed readiness", "pass-rate and failures roll up into Release Intelligence"),
    ],
    notes="The runner closes the loop — execute, stream logs, auto-heal failures, and feed outcomes into the release picture.",
)

# 10. Localization & accessibility
content_slide(
    "Localization & Accessibility",
    "Quality beyond functionality",
    [
        ("Localization Tester", "captures the page and AI-analyzes translation, missing keys & layout overflow"),
        ("Accessibility Scanner", "axe-core checks + AI audit across WCAG severities"),
        ("Actionable output", "scored results with issues you can log to Jira"),
        ("Project-scoped history", "results tracked per project over time"),
    ],
    notes="AAQUA also covers non-functional quality: localization and accessibility, both AI-augmented and tracked per project.",
)

# 11. Security engine
content_slide(
    "AI Secure Engine",
    "Security testing built in",
    [
        ("OWASP ZAP scanning", "baseline (passive), active (attack) and API (spec-import) scans"),
        ("AI vulnerability analysis", "LLM-assisted triage and explanation of findings"),
        ("SSRF-protected targets", "scan targets validated; private ranges gated by policy"),
        ("Jira integration", "raise defects from findings directly"),
    ],
    notes="A dedicated security subsystem wraps OWASP ZAP with AI triage and ties into governance and Jira.",
)

# 12. Release intelligence & governance
content_slide(
    "Release Intelligence & Governance",
    "A single signal for 'are we ready to ship?'",
    [
        ("Aggregated readiness", "combines automation, accessibility, localization & security results per project"),
        ("Release gating", "blocks release when critical + high findings exceed policy thresholds"),
        ("Project-scoped workspaces", "every activity bound to a project & its target app"),
        ("Dashboards", "trends and current quality posture at a glance"),
    ],
    notes="Governance turns all the module outputs into one release-readiness signal, with automated gating against thresholds.",
)

# 13. Integrations & identity
content_slide(
    "Integrations, Identity & Access",
    "Enterprise-ready by design",
    [
        ("Keycloak (OIDC)", "single sign-on, code+PKCE flow, role-based access control"),
        ("Jira", "log defects/issues straight from results"),
        ("On-prem LLM", "OpenAI-compatible local model — data stays internal"),
        ("Shared-infra deployment", "containerized, path-prefixed multi-tenant hosting"),
    ],
    notes="Identity is delegated to Keycloak (no local passwords); Jira for defect flow; on-prem LLM for privacy; deployed on shared container infra.",
)

# 14. Tech stack
content_slide(
    "Technology & Deployment",
    "Modern, containerized stack",
    [
        ("Frontend", "React 19 + React Router + Vite"),
        ("Backend", "Node.js / Express 5"),
        ("AI", "on-prem LLM (gpt-oss-20b), OpenAI-compatible API"),
        ("Engines", "OWASP ZAP, Playwright, axe-core"),
        ("Data & auth", "PostgreSQL (Sequelize ORM), Keycloak"),
        ("Delivery", "Docker / shared-infra, Nginx path-prefix routing"),
    ],
    two_col=True,
    notes="The stack slide for technical audiences.",
)

# 15. Summary
content_slide(
    "AAQUA in a Sentence",
    "",
    [
        ("One AI-powered platform", "for the entire QA lifecycle — generate, run, secure, and ship with confidence"),
        ("Breadth", "functional, API, security, localization, accessibility & release governance"),
        ("Trust", "on-prem AI, enterprise auth, and built-in security"),
        ("Outcome", "less manual effort, consistent quality, and a clear release signal"),
    ],
    notes="Close on the one-liner and the four pillars: breadth, AI, trust, outcomes.",
)

out = r"D:\AITesting\docs\AAQUA-Overview.pptx"
prs.save(out)
print("Saved:", out, "with", len(prs.slides._sldIdLst), "slides")
